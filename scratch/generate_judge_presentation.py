import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 Widescreen (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    BG_DARK = RGBColor(15, 23, 42)       # Slate 900
    BG_LIGHT = RGBColor(248, 250, 252)   # Slate 50
    CARD_BG = RGBColor(255, 255, 255)    # Pure White
    CARD_BORDER = RGBColor(226, 232, 240) # Slate 200
    
    PRIMARY_INDIGO = RGBColor(79, 70, 229)  # Indigo 600
    ACCENT_TEAL = RGBColor(13, 148, 136)   # Teal 600
    TEXT_MAIN = RGBColor(15, 23, 42)      # Slate 900
    TEXT_MUTED = RGBColor(100, 116, 139)  # Slate 500
    TEXT_LIGHT = RGBColor(241, 245, 249)  # Slate 100
    TEXT_SUB = RGBColor(148, 163, 184)    # Slate 400

    SEV_ROSE = RGBColor(225, 29, 72)      # Rose 600
    SEV_AMBER = RGBColor(217, 119, 6)     # Amber 600
    SEV_EMERALD = RGBColor(16, 185, 129)  # Emerald 500

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, tag, title, subtitle=None, is_dark=False):
        # Category Tag
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.35))
        tf = tag_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = tag.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_INDIGO if not is_dark else RGBColor(129, 140, 248)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.7))
        tf2 = title_box.text_frame
        tf2.word_wrap = True
        tf2.margin_left = tf2.margin_top = tf2.margin_right = tf2.margin_bottom = 0
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_LIGHT if is_dark else TEXT_MAIN

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.45))
            tf3 = sub_box.text_frame
            tf3.word_wrap = True
            tf3.margin_left = tf3.margin_top = tf3.margin_right = tf3.margin_bottom = 0
            p3 = tf3.paragraphs[0]
            p3.text = subtitle
            p3.font.size = Pt(13)
            p3.font.color.rgb = TEXT_SUB if is_dark else TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def set_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # =========================================================================
    # SLIDE 1: COVER
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BG_DARK
    bg1.line.fill.background()

    # Brand Title
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "SHIKSHASETU  •  शिक्षासेतु"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(129, 140, 248) # Indigo 400
    p.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "The School ERP That Actually\nUnderstands Learning."
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "Connecting teachers, students, and parents around one learner.\nOne canonical dataset. Three useful real-time experiences."
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(148, 163, 184)

    set_notes(s1, "SPEAKER NOTES (15s):\nGood morning judges. Traditional school ERPs are great at digitizing school records, but they stop at dashboards. ShikshaSetu is built on a fundamentally different principle: turning everyday school information into coordinated next actions for teachers, students, and parents around one learner.")

    # =========================================================================
    # SLIDE 2: THE PROBLEM
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "The Real-World Problem", "Schools Have Plenty of Data. But Data Doesn't Automatically Become Action.", "Schools generate millions of data points every term, yet information remains disconnected from daily action.")

    # 3 Column Cards
    col_w = Inches(3.64)
    gap = Inches(0.39)
    top_pos = Inches(2.2)
    card_h = Inches(4.5)

    # Col 1: Teacher Question
    add_card(s2, Inches(0.8), top_pos, col_w, card_h)
    tb = s2.shapes.add_textbox(Inches(1.1), top_pos + Inches(0.3), col_w - Inches(0.6), card_h - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "👨‍🏫 THE TEACHER"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_INDIGO
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\"Who needs my attention today?\""
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "Teachers manage 40+ students per class. Marks are entered into spreadsheets, but finding which 3 students are falling behind in a specific topic requires manual cross-referencing."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # Col 2: Student Question
    add_card(s2, Inches(0.8) + col_w + gap, top_pos, col_w, card_h)
    tb = s2.shapes.add_textbox(Inches(1.1) + col_w + gap, top_pos + Inches(0.3), col_w - Inches(0.6), card_h - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎓 THE STUDENT"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\"What should I study next?\""
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "Students see a raw score like 58% in Mathematics, but the portal doesn't explain the underlying gap (Equivalent Fractions) or give them a 5-minute practice loop."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    # Col 3: Parent Question
    add_card(s2, Inches(0.8) + (col_w + gap)*2, top_pos, col_w, card_h)
    tb = s2.shapes.add_textbox(Inches(1.1) + (col_w + gap)*2, top_pos + Inches(0.3), col_w - Inches(0.6), card_h - Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏡 THE PARENT"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SEV_AMBER
    p.space_after = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "\"How can I help my child?\""
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "Parents receive report cards months after exams. They don't know what questions to ask at dinner to reinforce what the teacher taught in class today."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MUTED

    set_notes(s2, "SPEAKER NOTES (30s):\nSchools generate huge amounts of data: attendance, marks, homework, and behavior logs. But the data remains trapped in administrative silos. The teacher struggles to know who needs attention right now, the student doesn't know what to study next, and the parent doesn't know how to support at home.")

    # =========================================================================
    # SLIDE 3: THE PARADIGM SHIFT
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Paradigm Shift", "From Recording What Happened to Deciding What Happens Next", "We do not replace school data — we turn passive records into active, connected workflows.")

    # Left Box: Traditional ERP
    add_card(s3, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5), bg_color=RGBColor(241, 245, 249))
    tb = s3.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TRADITIONAL SCHOOL ERP"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(16)

    flow1 = [
        "1. Attendance Taken → Logged into Database",
        "2. Marks Entered → Stored in Exam Table",
        "3. Homework Assigned → Saved in Portal",
        "4. End of Term → Generate PDF Summary"
    ]
    for step in flow1:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "❌ Result: Passive archive of past events."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SEV_ROSE

    # Right Box: ShikshaSetu
    add_card(s3, Inches(6.9), Inches(2.2), Inches(5.6), Inches(4.5), bg_color=CARD_BG, border_color=PRIMARY_INDIGO)
    tb = s3.shapes.add_textbox(Inches(7.2), Inches(2.5), Inches(5.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "SHIKSHASETU CONNECTED ECOSYSTEM"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_INDIGO
    p.space_after = Pt(16)

    flow2 = [
        "1. Teacher logs 58% in Equivalent Fractions",
        "2. Teacher Radar: Suggests 10-min visual review",
        "3. Student Portal: Prepares 5-min practice & diagnostic quiz",
        "4. Parent Portal: Generates practical dinner question prompt"
    ]
    for step in flow2:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MAIN
        p.space_after = Pt(10)

    p = tf.add_paragraph()
    p.text = "✓ Result: Continuous learning loop and coordinated action."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = SEV_EMERALD

    set_notes(s3, "SPEAKER NOTES (25s):\nTraditional ERPs stop at 'What happened?'. ShikshaSetu asks: 'Why does it matter?' and 'What should happen next?'. When a mark or attendance record is entered, our rules engine and contextual layer immediately generate role-specific actions for the teacher, student, and parent.")

    # =========================================================================
    # SLIDE 4: THE CORE ARCHITECTURE
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "The Single Source of Truth", "One Student. One Connected Source of Truth. Three Useful Experiences.", "Data is entered once by teachers and administrators. Every portal reads the exact same canonical record.")

    # 4 Quadrants / Hub Layout
    add_card(s4, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    
    # Center Hub
    add_card(s4, Inches(4.8), Inches(3.2), Inches(3.7), Inches(2.2), bg_color=PRIMARY_INDIGO, border_color=None)
    tb = s4.shapes.add_textbox(Inches(5.0), Inches(3.5), Inches(3.3), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CANONICAL LEARNER"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(224, 231, 255)
    p.space_after = Pt(4)
    
    p2 = tf.add_paragraph()
    p2.text = "Aarav Sharma · 8A\nMath 58% (Fractions)"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)

    # Top Left: Teacher
    add_card(s4, Inches(1.1), Inches(2.5), Inches(3.3), Inches(1.6), bg_color=RGBColor(248, 250, 252))
    tb = s4.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(3.1), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "👩‍🏫 TEACHER WORKSPACE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_INDIGO
    p2 = tf.add_paragraph()
    p2.text = "• 3-Student Attention Radar\n• Instant Class Attendance\n• 1-Click CBSE Report Cards"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MAIN

    # Top Right: Student
    add_card(s4, Inches(8.9), Inches(2.5), Inches(3.3), Inches(1.6), bg_color=RGBColor(248, 250, 252))
    tb = s4.shapes.add_textbox(Inches(9.0), Inches(2.6), Inches(3.1), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🎓 STUDENT COMMAND CENTER"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p2 = tf.add_paragraph()
    p2.text = "• Today's Action Priorities\n• AI Diagnostic Revision Notes\n• SchoolMitra Homework Tutor"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MAIN

    # Bottom Right: Parent
    add_card(s4, Inches(8.9), Inches(4.7), Inches(3.3), Inches(1.6), bg_color=RGBColor(248, 250, 252))
    tb = s4.shapes.add_textbox(Inches(9.0), Inches(4.8), Inches(3.1), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏡 PARENT COMPANION"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SEV_AMBER
    p2 = tf.add_paragraph()
    p2.text = "• 3-Subject Progress Snapshot\n• 5-Min Home Dinner Prompts\n• Live GPS School Bus Marker"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MAIN

    # Bottom Left: Operations
    add_card(s4, Inches(1.1), Inches(4.7), Inches(3.3), Inches(1.6), bg_color=RGBColor(248, 250, 252))
    tb = s4.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(3.1), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🏛️ ADMIN & OPERATIONS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p2 = tf.add_paragraph()
    p2.text = "• Operational Needs Queue\n• Master Registry CRUD\n• Fast Camera Gate QR Pass"
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MAIN

    set_notes(s4, "SPEAKER NOTES (20s):\nOur architecture enforces strict canonical single-source-of-truth. We never duplicate records. Data is stored in normalized Supabase PostgreSQL with Row Level Security, and each portal views the same data customized for their role.")

    # =========================================================================
    # SLIDE 5: REAL STUDENT JOURNEY (AARAV SHARMA)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Real MVP Data Journey", "Same Data. Different Action.", "Tracing canonical student Aarav Sharma (Class 8A) across all three portals.")

    # Left Column: The Fact
    add_card(s5, Inches(0.8), Inches(2.2), Inches(3.4), Inches(4.5), bg_color=RGBColor(254, 242, 242), border_color=RGBColor(254, 205, 211))
    tb = s5.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(3.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "CANONICAL RECORD"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = SEV_ROSE
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "Aarav Sharma\nClass 8A · Roll #801"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(16)

    p = tf.add_paragraph()
    p.text = "📊 Math Score: 58%\n🔍 Gap: Equivalent Fractions\n🔬 Science Score: 82%\n📖 English Score: 76%\n📅 Attendance: 94%"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    # Right Column: 3 Real Actions
    add_card(s5, Inches(4.5), Inches(2.2), Inches(8.0), Inches(4.5))
    tb = s5.shapes.add_textbox(Inches(4.8), Inches(2.4), Inches(7.4), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True

    # 1. Teacher Action
    p = tf.paragraphs[0]
    p.text = "1. TEACHER PORTAL (/teacher) — Ms. Ananya Mehra"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_INDIGO
    p = tf.add_paragraph()
    p.text = "→ Radar highlights Aarav in top 3 attention list with 58% in Mathematics.\n→ AI Suggestion: \"Run a 10-minute visual fraction strips review before Friday's exam.\""
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)

    # 2. Student Action
    p = tf.add_paragraph()
    p.text = "2. STUDENT PORTAL (/student) — Aarav Sharma"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p = tf.add_paragraph()
    p.text = "→ Action Center displays priority: 🟠 Friday Test: Equivalent Fractions.\n→ AI Revision Notes: Explains pizza analogy, gives 3-question diagnostic quiz with instant feedback."
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MAIN
    p.space_after = Pt(12)

    # 3. Parent Action
    p = tf.add_paragraph()
    p.text = "3. PARENT PORTAL (/parent) — Sunita Sharma"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = SEV_AMBER
    p = tf.add_paragraph()
    p.text = "→ Progress card clearly notes: \"Mathematics needs a little attention this week (58%).\"\n→ 5-Min Dinner Prompt: \"Ask Aarav: Can you explain why 2/4 of a pizza is the same as 1/2?\""
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MAIN

    set_notes(s5, "SPEAKER NOTES (30s):\nLet's look at a concrete example from our codebase: Aarav Sharma in Class 8A. When his 58% score in Equivalent Fractions is recorded, Ms. Ananya gets a teaching recommendation, Aarav gets a 5-minute revision quiz on that exact gap, and his mother Sunita gets a practical question to ask at dinner.")

    # =========================================================================
    # SLIDE 6: THE LEARNING LOOP
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "The Continuous Improvement Cycle", "Data Becomes a Continuous Learning Loop", "Every assessment feeds back into the student's mastery trajectory.")

    steps = [
        ("1. OBSERVE", "Teacher records test marks & daily attendance.", PRIMARY_INDIGO),
        ("2. UNDERSTAND", "Rules engine identifies topic-level learning gap.", ACCENT_TEAL),
        ("3. ACT", "Portal prioritizes Friday exam prep & homework.", SEV_AMBER),
        ("4. LEARN", "NCERT revision note with interactive concept chips.", PRIMARY_INDIGO),
        ("5. CHECK", "3-Question diagnostic quiz evaluates mastery level.", SEV_EMERALD),
        ("6. REINFORCE", "Parent reinforces concept at home with guided prompt.", PRIMARY_INDIGO)
    ]

    for i, (title, desc, color) in enumerate(steps):
        row = i // 3
        col = i % 3
        card_x = Inches(0.8) + col * Inches(3.95)
        card_y = Inches(2.3) + row * Inches(2.2)
        
        add_card(s6, card_x, card_y, Inches(3.7), Inches(2.0))
        tb = s6.shapes.add_textbox(card_x + Inches(0.2), card_y + Inches(0.25), Inches(3.3), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MAIN

    set_notes(s6, "SPEAKER NOTES (20s):\nThis creates a closed learning loop: Observe, Understand, Act, Learn, Check, and Reinforce. The school and family operate as a single coordinated support team around the student.")

    # =========================================================================
    # SLIDE 7: AI WITH CONTEXT (NOT A CHATBOT)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Contextual Intelligence", "AI With Context. Not Just a Chatbot.", "AI is not the product — it is an intelligent contextual layer over authorized school facts.")

    add_card(s7, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tb = s7.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "HOW OUR CONTEXT PIPELINE WORKS:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_INDIGO
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "Student Academic Gaps + Role Boundaries + NCERT Curriculum Benchmarks  →  Google Gemini API / Rule Fallbacks  →  Actionable Output"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(16)

    examples = [
        ("🎓 SchoolMitra (Student AI)", "Curriculum-grounded Socratic tutor that explains difficult NCERT concepts step-by-step without doing the homework for them."),
        ("👨‍🏫 Teacher Copilot", "Generates lesson plans, diagnostic review suggestions, and automated CBSE-standard report card remarks based on actual student evidence."),
        ("🏡 Parent AI Guide", "Answers parental questions in everyday language while strictly respecting student privacy and role boundaries.")
    ]

    for title, desc in examples:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    set_notes(s7, "SPEAKER NOTES (25s):\nJudges often ask: 'Isn't this just ChatGPT inside an ERP?'. The answer is no. AI in ShikshaSetu is strictly grounded in verified school context: syllabus benchmarks, actual student marks, and role boundaries. It never sees unauthorized data and never fabricates scores.")

    # =========================================================================
    # SLIDE 8: DATA INTEGRITY & WHO ENTERS THE DATA
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Data Provenance", "Who Enters the Data? Data is Entered Once at the Source.", "No synthetic metrics, no duplicate student IDs, and no fake demo data.")

    actors = [
        ("👨‍🏫 Teacher", "Logs attendance, marks exam scores, assigns homework, and records classroom observations.", PRIMARY_INDIGO),
        ("🏛️ Administrator", "Maintains canonical student rosters, teacher profiles, class allocations, and fee records.", TEXT_MAIN),
        ("🚌 Bus Driver", "Broadcasts live hardware device GPS coordinates via native browser geolocation.", SEV_AMBER),
        ("🎓 Student", "Completes revision notes, takes diagnostic concept checks, and logs emotional check-ins.", ACCENT_TEAL),
        ("🏡 Parent", "Requests digital gate passes and initiates direct teacher communications.", SEV_EMERALD)
    ]

    for i, (actor, desc, color) in enumerate(actors):
        card_y = Inches(2.2) + i * Inches(0.9)
        add_card(s8, Inches(0.8), card_y, Inches(11.7), Inches(0.78))
        tb = s8.shapes.add_textbox(Inches(1.0), card_y + Inches(0.12), Inches(11.3), Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{actor}:  "
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color
        
        # Add desc inline
        p2 = tf.paragraphs[0].add_run()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.bold = False
        p2.font.color.rgb = TEXT_MAIN

    set_notes(s8, "SPEAKER NOTES (20s):\nWho creates the data? Real stakeholders at the source. Teachers enter marks, admins manage classes, drivers stream GPS, and parents request gate passes. AI does not invent facts; it turns these verified human inputs into coordinated action.")

    # =========================================================================
    # SLIDE 9: WHAT WE ACTUALLY BUILT (MVP SCOPE)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Functional Implementation", "What We Actually Built: 6 Operational Portals", "Everything shown is verified, tested, and running in production.")

    categories = [
        ("🎓 LEARNING", ["• Daily Priorities (Due Today / Tests)", "• AI Revision Notes + Concept Chips", "• 3-Question Diagnostic Quiz Loop", "• SchoolMitra NCERT Doubt Tutor"]),
        ("👨‍🏫 TEACHING", ["• 3-Student Attention Radar", "• Rapid 1-Tap Attendance Taker", "• Academic Marks Panel", "• CBSE Printable PDF Report Cards"]),
        ("🏡 FAMILY", ["• Academic Progress Snapshot", "• 5-Minute Dinner Action Prompts", "• Digital Gate Pass QR Generator", "• Teacher Direct Messaging Thread"]),
        ("🚌 OPERATIONS", ["• Driver Live GPS Console (watchPosition)", "• Leaflet OpenStreetMap Tracking", "• Gate Camera QR Scanner", "• Admin Master Registry CRUD"])
    ]

    for i, (cat, items) in enumerate(categories):
        col = i % 4
        card_x = Inches(0.8) + col * Inches(2.98)
        add_card(s9, card_x, Inches(2.2), Inches(2.8), Inches(4.5))
        tb = s9.shapes.add_textbox(card_x + Inches(0.15), Inches(2.4), Inches(2.5), Inches(4.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cat
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_INDIGO
        p.space_after = Pt(12)

        for it in items:
            p = tf.add_paragraph()
            p.text = it
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT_MAIN
            p.space_after = Pt(6)

    set_notes(s9, "SPEAKER NOTES (25s):\nAcross 6 functional portals, we have built a complete learning and operations ecosystem: student command center, teacher radar, parent companion, driver GPS console, gate pass verification, and admin master registry.")

    # =========================================================================
    # SLIDE 10: LIVE MVP DEMO FLOW
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    bg10 = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg10.fill.solid()
    bg10.fill.fore_color.rgb = BG_DARK
    bg10.line.fill.background()

    tb = s10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "LIVE MVP DEMONSTRATION"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(129, 140, 248)
    p.space_after = Pt(14)

    p = tf.add_paragraph()
    p.text = "Let's Follow One Student: Aarav Sharma"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(24)

    demo_steps = [
        "1. Teacher Portal  →  Ms. Ananya sees Aarav highlighted with 58% in Equivalent Fractions.",
        "2. Student Portal  →  Aarav opens Friday test prep and takes a 3-question diagnostic revision quiz.",
        "3. Parent Portal   →  Sunita sees the same 58% and gets a 5-minute dinner question prompt.",
        "4. Driver Portal   →  Driver Rajesh starts trip with real GPS; Parent watches bus move on live map."
    ]

    for st in demo_steps:
        p = tf.add_paragraph()
        p.text = st
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(10)

    set_notes(s10, "SPEAKER NOTES (15s):\nNow, let's switch directly to the live application and follow Aarav Sharma through the Teacher, Student, Parent, and Driver portals.")

    # =========================================================================
    # SLIDE 11: CONNECTED OPERATIONS — LIVE GPS TRACKING
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "Connected School Operations", "Real Hardware Device GPS Telemetry Pipeline", "True device geolocation streaming without simulated mock movement.")

    add_card(s11, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tb = s11.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "THE 4-STAGE TELEMETRY PIPELINE:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_INDIGO
    p.space_after = Pt(16)

    pipe = [
        ("📱 1. Driver Browser GPS", "Acquires real device coordinates via navigator.geolocation.watchPosition ({ enableHighAccuracy: true })."),
        ("🛡️ 2. Server Authorization", "Server Action validates driver identity, updates canonical buses table, and records precision telemetry."),
        ("⚡ 3. Supabase Realtime", "Broadcasts coordinates over school:sch-demo-001:bus:BUS-21 channel at 4-second intervals."),
        ("🗺️ 4. Parent Portal Map", "Leaflet OpenStreetMap animates bus location marker with real freshness badge (e.g., \"Updated 4s ago\").")
    ]

    for title, desc in pipe:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(8)

    set_notes(s11, "SPEAKER NOTES (20s):\nOur live bus tracking uses actual browser hardware GPS from the driver's phone. Coordinates are written to PostgreSQL and broadcast via WebSockets directly to the parent's map.")

    # =========================================================================
    # SLIDE 12: CONCLUSION
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    bg12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg12.fill.solid()
    bg12.fill.fore_color.rgb = BG_DARK
    bg12.line.fill.background()

    tb = s12.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "SHIKSHASETU  •  SUMMARY"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(129, 140, 248)
    p.space_after = Pt(16)

    p2 = tf.add_paragraph()
    p2.text = "Don't Just Manage Your School.\nUnderstand Every Learner."
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.space_after = Pt(20)

    p3 = tf.add_paragraph()
    p3.text = "DATA  →  UNDERSTANDING  →  ACTION  →  OUTCOME\n\nThank you. We welcome your questions."
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(226, 232, 240)

    set_notes(s12, "SPEAKER NOTES (15s):\nIn summary, ShikshaSetu connects school data into immediate action for every teacher, student, and parent. Thank you, and we welcome your questions.")

    # =========================================================================
    # APPENDIX SLIDES
    # =========================================================================
    
    # Appendix 1: Architecture & Stack
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Appendix: Architecture", "Full-Stack System Architecture & Technology Stack", "Production-ready enterprise Next.js and Supabase infrastructure.")
    add_card(s13, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tb = s13.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Frontend Framework: Next.js 14 (App Router, React 18, TypeScript, Tailwind CSS, Framer Motion)\n• Database: Supabase PostgreSQL with 38 modular SQL migrations & Row Level Security (RLS)\n• Authentication: Clerk Auth with idempotent server-side onboarding linking to canonical tables\n• Real-Time Protocol: Supabase Realtime (PostgreSQL WAL replication) + Socket.io fallback\n• Mapping & Telemetry: Leaflet.js, OpenStreetMap, HTML5 Geolocation API (watchPosition)\n• Quality Assurance: Vitest (250 tests passed across 36 test files), Playwright E2E"
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MAIN
    set_notes(s13, "APPENDIX: Detailed stack specifications for technical jury inquiries.")

    # Appendix 2: Security & Tenant Isolation
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "Appendix: Security", "Multi-Tenant Isolation & Row Level Security (RLS)", "Strict data boundary enforcement across school tenants.")
    add_card(s14, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    tb = s14.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Tenant Scoping: Every table includes school_id partition with enforced RLS policies.\n• Guardian Boundaries: Parents can strictly query only students linked via guardian_access table.\n• Teacher Boundaries: Teachers access student records within their assigned grade/section.\n• Sensitive Notes Privacy: Teacher private notes and medical flags are masked from parent feeds.\n• Server Action Authorization: All mutations verify authenticated role context before execution."
    p.font.size = Pt(12)
    p.font.color.rgb = TEXT_MAIN
    set_notes(s14, "APPENDIX: Security and RLS architecture details.")

    # Save presentation to disk
    output_path = os.path.join(os.getcwd(), "ShikshaSetu_Judge_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
